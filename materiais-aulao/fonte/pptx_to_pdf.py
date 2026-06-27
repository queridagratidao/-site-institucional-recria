# -*- coding: utf-8 -*-
# Conversor genérico .pptx -> .pdf (sem LibreOffice), aproximado: redesenha
# retangulos/textos/imagens na mesma posicao. Formas especiais (chevron, oval)
# saem como caixa retangular equivalente. Serve pra visualizar/compartilhar;
# pra fidelidade total, abrir o .pptx no PowerPoint/Google Slides e exportar PDF.
import sys
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_SHAPE
from reportlab.pdfgen import canvas
from reportlab.lib.colors import Color
from reportlab.lib.utils import ImageReader
import io

EMU_PER_INCH = 914400
PT_PER_INCH = 72

def emu_to_pt(emu):
    return emu / EMU_PER_INCH * PT_PER_INCH

def rgb_to_color(rgb):
    if rgb is None:
        return None
    return Color(rgb[0] / 255, rgb[1] / 255, rgb[2] / 255)

def get_fill_color(shape):
    try:
        if shape.fill.type is not None and shape.fill.type == 1:  # solid
            c = shape.fill.fore_color
            if c.type is not None:
                rgb = c.rgb
                return rgb_to_color((rgb[0], rgb[1], rgb[2]))
    except Exception:
        pass
    return None

def get_line(shape):
    try:
        line = shape.line
        if line.fill.type == 1:
            rgb = line.color.rgb
            width_pt = line.width.pt if line.width else 1
            return rgb_to_color((rgb[0], rgb[1], rgb[2])), width_pt
    except Exception:
        pass
    return None, None

def is_oval(shape):
    try:
        return shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and shape.auto_shape_type == MSO_SHAPE.OVAL
    except Exception:
        return False

def convert(pptx_path, pdf_path):
    prs = Presentation(pptx_path)
    pw_pt = emu_to_pt(prs.slide_width)
    ph_pt = emu_to_pt(prs.slide_height)
    c = canvas.Canvas(pdf_path, pagesize=(pw_pt, ph_pt))

    for slide in prs.slides:
        # background
        bg_color = None
        try:
            bgfill = slide.background.fill
            if bgfill.type == 1:
                rgb = bgfill.fore_color.rgb
                bg_color = rgb_to_color((rgb[0], rgb[1], rgb[2]))
        except Exception:
            pass
        if bg_color:
            c.setFillColor(bg_color)
            c.rect(0, 0, pw_pt, ph_pt, fill=1, stroke=0)

        for shape in slide.shapes:
            x = emu_to_pt(shape.left) if shape.left is not None else 0
            y = emu_to_pt(shape.top) if shape.top is not None else 0
            w = emu_to_pt(shape.width) if shape.width is not None else 0
            h = emu_to_pt(shape.height) if shape.height is not None else 0
            y_pdf = ph_pt - y - h  # PDF origin é embaixo

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    img_bytes = shape.image.blob
                    img_reader = ImageReader(io.BytesIO(img_bytes))
                    c.drawImage(img_reader, x, y_pdf, width=w, height=h,
                                preserveAspectRatio=True, mask='auto')
                except Exception as e:
                    print("falha ao desenhar imagem:", e)
                continue

            if shape.shape_type is not None and shape.shape_type != 17:  # nao e textbox puro
                fill_color = get_fill_color(shape)
                line_color, line_width = get_line(shape)
                do_fill = 1 if fill_color else 0
                do_stroke = 1 if line_color else 0
                if do_fill or do_stroke:
                    if do_fill:
                        c.setFillColor(fill_color)
                    if do_stroke:
                        c.setStrokeColor(line_color)
                        c.setLineWidth(line_width or 1)
                    if is_oval(shape):
                        c.ellipse(x, y_pdf, x + w, y_pdf + h, fill=do_fill, stroke=do_stroke)
                    else:
                        c.roundRect(x, y_pdf, w, h, 6, fill=do_fill, stroke=do_stroke)

            if shape.has_text_frame:
                tf = shape.text_frame
                max_width = w - 8 if w else 400
                # monta lista de (texto, tamanho, fonte, cor) por linha, respeitando \n manuais
                render_lines = []
                for para in tf.paragraphs:
                    runs = para.runs
                    align_name = para.alignment.name if para.alignment is not None else None
                    if not runs:
                        render_lines.append(("", 14, "Helvetica", Color(0, 0, 0), align_name))
                        continue
                    full_text = "".join(r.text for r in runs)
                    r0 = runs[0]
                    size = r0.font.size.pt if r0.font.size else 14
                    bold = r0.font.bold or False
                    color = rgb_to_color((r0.font.color.rgb[0], r0.font.color.rgb[1], r0.font.color.rgb[2])) \
                        if r0.font.color and r0.font.color.type is not None else Color(0, 0, 0)
                    font_name = "Helvetica-Bold" if bold else "Helvetica"
                    for manual_line in full_text.split("\n"):
                        words = manual_line.split(" ")
                        cur = ""
                        for word in words:
                            test = (cur + " " + word).strip()
                            if c.stringWidth(test, font_name, size) <= max_width or not cur:
                                cur = test
                            else:
                                render_lines.append((cur, size, font_name, color, align_name))
                                cur = word
                        render_lines.append((cur, size, font_name, color, align_name))

                line_h_list = [size * 1.25 for (_, size, _, _, _) in render_lines]
                total_h = sum(line_h_list) + 8
                anchor = tf.vertical_anchor
                if anchor is not None and anchor.name == "MIDDLE":
                    start_y = y + max(0, (h - total_h) / 2)
                elif anchor is not None and anchor.name == "BOTTOM":
                    start_y = y + max(0, h - total_h)
                else:
                    start_y = y + 4

                line_y = start_y
                for text, size, font_name, color, align_name in render_lines:
                    c.setFont(font_name, size)
                    c.setFillColor(color)
                    ty = ph_pt - line_y - size
                    lw = c.stringWidth(text, font_name, size)
                    if align_name == "CENTER":
                        lx = x + (w - lw) / 2 if w else x
                    elif align_name == "RIGHT":
                        lx = x + w - lw - 4 if w else x
                    else:
                        lx = x + 4
                    c.drawString(lx, ty, text)
                    line_y += size * 1.25

        c.showPage()

    c.save()
    print("PDF salvo:", pdf_path)

if __name__ == "__main__":
    convert(sys.argv[1], sys.argv[2])
