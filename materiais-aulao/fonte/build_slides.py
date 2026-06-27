# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

BLACK = RGBColor(0x0A, 0x0A, 0x0A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GOLD = RGBColor(0xC9, 0xA2, 0x4E)
GOLD_SOFT = RGBColor(0xE8, 0xD3, 0x9A)
GRAY = RGBColor(0x6B, 0x6B, 0x6B)
CARD_BG = RGBColor(0x16, 0x16, 0x16)

TITLE_FONT = "Cambria"
BODY_FONT = "Calibri"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

def add_slide(bg=BLACK):
    s = prs.slides.add_slide(BLANK)
    bgfill = s.background.fill
    bgfill.solid()
    bgfill.fore_color.rgb = bg
    return s

def textbox(slide, x, y, w, h, text, size=18, color=WHITE, bold=False, italic=False,
            font=BODY_FONT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.0,
            space_after=0):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
        r.font.italic = italic
        r.font.name = font
    return tb

def bullets(slide, x, y, w, h, items, size=15, color=WHITE, font=BODY_FONT,
            space_after=10, line_spacing=1.08, bullet_color=GOLD):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(space_after)
        p.line_spacing = line_spacing
        pPr = p._p.get_or_add_pPr()
        buChar = pPr.makeelement(qn('a:buChar'), {'char': '—'})
        buFont = pPr.makeelement(qn('a:buFont'), {'typeface': font})
        buClr = pPr.makeelement(qn('a:buClr'))
        srgb = buClr.makeelement(qn('a:srgbClr'), {'val': '%02X%02X%02X' % (bullet_color[0], bullet_color[1], bullet_color[2])})
        buClr.append(srgb)
        pPr.append(buClr)
        pPr.append(buFont)
        pPr.append(buChar)
        pPr.set('marL', '228600')
        pPr.set('indent', '-228600')
        if isinstance(item, tuple):
            head, body = item
            r1 = p.add_run()
            r1.text = head + "  "
            r1.font.bold = True
            r1.font.size = Pt(size)
            r1.font.color.rgb = color
            r1.font.name = font
            r2 = p.add_run()
            r2.text = body
            r2.font.size = Pt(size)
            r2.font.color.rgb = color
            r2.font.name = font
        else:
            r = p.add_run()
            r.text = item
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.name = font
    return tb

def rect(slide, x, y, w, h, fill=CARD_BG, line=None, radius=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    sp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    if radius:
        sp.adjustments[0] = 0.06
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line:
        sp.line.color.rgb = line
        sp.line.width = Pt(0.75)
    else:
        sp.line.fill.background()
    sp.shadow.inherit = False
    return sp

def oval(slide, x, y, w, h, fill=GOLD):
    sp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill; sp.line.fill.background(); sp.shadow.inherit = False
    return sp

def kicker(slide, text, x=0.9, y=0.6, color=GOLD):
    return textbox(slide, x, y, 9.0, 0.4, text.upper(), size=12, color=color, bold=True, font=BODY_FONT)

def pagenum(slide, n):
    textbox(slide, 12.5, 7.05, 0.6, 0.35, str(n), size=11, color=GRAY, font=BODY_FONT, align=PP_ALIGN.RIGHT)

LOGO_ICON = "assets/logo_icon.png"

def logo_mark(slide, x=0.9, y=0.55):
    slide.shapes.add_picture(LOGO_ICON, Inches(x), Inches(y - 0.03), height=Inches(0.32))
    textbox(slide, x + 0.4, y - 0.04, 2.5, 0.3, "RECRIA", size=12, color=GOLD, bold=True, font=BODY_FONT)

n = 0
def nextn():
    global n
    n += 1
    return n

def block_title(num, title, sub):
    s = add_slide(BLACK)
    textbox(s, 0.9, 2.9, 2.0, 1.2, f"{num:02d}", size=80, color=GOLD, bold=True, font=TITLE_FONT)
    textbox(s, 3.2, 3.1, 9.0, 1.1, title, size=30, color=WHITE, bold=True, font=TITLE_FONT, line_spacing=1.05)
    textbox(s, 3.2, 4.15, 8.6, 0.8, sub, size=15, color=GOLD_SOFT, italic=True, font=BODY_FONT)
    pagenum(s, nextn())
    return s

def quote_slide(text, sub=""):
    s = add_slide(GOLD)
    textbox(s, 1.4, 2.3, 10.5, 2.6, text, size=27, color=BLACK, bold=True, font=TITLE_FONT, line_spacing=1.25)
    if sub:
        textbox(s, 1.4, 5.1, 10.0, 0.6, sub, size=15, color=BLACK, italic=True, font=BODY_FONT)
    pagenum(s, nextn())
    return s

def circle_flow(slide, items, centers_y, circ_d, spacing, startx, label_y, label_w,
                 num_size, label_size, gold_idxs=(), highlight_idx=None, highlight_box_h=2.95):
    n = len(items)
    centers_x = [startx + i * spacing for i in range(n)]

    # linha conectora continua atras dos circulos, garante alinhamento perfeito
    line_y = centers_y - 0.012
    line = rect(slide, centers_x[0], line_y, centers_x[-1] - centers_x[0], 0.024, fill=GOLD)
    line.line.fill.background()

    if highlight_idx is not None:
        hcx = centers_x[highlight_idx]
        hbox = rect(slide, hcx - label_w/2 - 0.25, centers_y - circ_d/2 - 0.3, label_w + 0.5, highlight_box_h, fill=BLACK, radius=True)
        hbox.fill.background()
        hbox.line.color.rgb = GOLD
        hbox.line.width = Pt(1.5)

    for i, (txt, cx) in enumerate(zip(items, centers_x)):
        is_gold = (i in gold_idxs) or (i == highlight_idx)
        fill = GOLD if is_gold else BLACK
        color = BLACK if is_gold else WHITE
        circ = oval(slide, cx - circ_d/2, centers_y - circ_d/2, circ_d, circ_d, fill=fill)
        circ.line.fill.solid()
        circ.line.color.rgb = GOLD
        circ.line.width = Pt(1.5)
        tf = circ.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = str(i + 1)
        r.font.size = Pt(num_size)
        r.font.bold = True
        r.font.color.rgb = color
        r.font.name = TITLE_FONT
        head, sub = txt if isinstance(txt, tuple) else (txt, None)
        textbox(slide, cx - label_w/2, label_y, label_w, 0.7, head, size=label_size,
                color=GOLD if is_gold else WHITE, bold=True, font=BODY_FONT,
                align=PP_ALIGN.CENTER, line_spacing=1.15)
        if sub:
            textbox(slide, cx - label_w/2, label_y + 0.55, label_w, 1.3, sub, size=label_size - 2.5,
                    color=GOLD_SOFT, font=BODY_FONT, align=PP_ALIGN.CENTER, line_spacing=1.2)
    return centers_x

def two_col_compare(titulo, colA_h, colA_items, colB_h, colB_items, colA_gold=False, colB_gold=True):
    s = add_slide(BLACK)
    textbox(s, 0.9, 1.0, 11, 0.8, titulo, size=27, color=WHITE, bold=True, font=TITLE_FONT, line_spacing=1.1)
    fillA = GOLD if colA_gold else CARD_BG
    fillB = GOLD if colB_gold else CARD_BG
    colorA = BLACK if colA_gold else WHITE
    colorB = BLACK if colB_gold else WHITE
    rect(s, 0.9, 2.1, 5.6, 4.5, fill=fillA, radius=True)
    textbox(s, 1.25, 2.45, 5.0, 0.5, colA_h, size=15, color=GOLD if not colA_gold else BLACK, bold=True, font=BODY_FONT)
    bullets(s, 1.25, 3.05, 4.9, 3.3, colA_items, size=14.5, color=colorA, bullet_color=colorA)
    rect(s, 6.85, 2.1, 5.6, 4.5, fill=fillB, radius=True)
    textbox(s, 7.2, 2.45, 5.0, 0.5, colB_h, size=15, color=GOLD if not colB_gold else BLACK, bold=True, font=BODY_FONT)
    bullets(s, 7.2, 3.05, 4.9, 3.3, colB_items, size=14.5, color=colorB, bullet_color=colorB)
    pagenum(s, nextn())
    return s

# ============================================================
# SLIDE 1: TITLE
# ============================================================
s = add_slide(BLACK)
logo_mark(s, 0.9, 0.6)
textbox(s, 0.9, 1.7, 11.5, 0.5, "AULA GRATUITA", size=16, color=GOLD, bold=True, font=BODY_FONT)
textbox(s, 0.9, 2.15, 11.3, 2.3, "Por que seu marketing não funciona\ne seu comercial não converte\ncomo deveria?", size=34, color=WHITE, bold=True, font=TITLE_FONT, line_spacing=1.08)
textbox(s, 0.9, 4.55, 11.0, 1.1, "Vou abrir o jogo com você: 10 anos de experiência em marketing e comercial, para te mostrar o que ninguém te mostrou até hoje. De graça, em uma única aula.", size=15.5, color=GOLD_SOFT, italic=True, font=BODY_FONT, line_spacing=1.25)
textbox(s, 0.9, 5.85, 11.0, 0.7, "Marketing e comercial alinhados podem triplicar as suas vendas. Nessa aula você vai descobrir como.", size=13.5, color=WHITE, font=BODY_FONT, line_spacing=1.25)
textbox(s, 0.9, 6.85, 6, 0.4, "por Amanda · Recria", size=13, color=GRAY, font=BODY_FONT)

# ============================================================
# SLIDE 2: QUEM SOU EU
# ============================================================
s = add_slide(BLACK)
kicker(s, "Antes de começar")
textbox(s, 0.9, 1.05, 8, 0.9, "Quem te ensina isso hoje", size=30, color=WHITE, bold=True, font=TITLE_FONT)
rect(s, 0.9, 2.1, 7.0, 4.4, fill=CARD_BG, radius=True)
bullets(s, 1.3, 2.5, 6.2, 3.6, [
    ("Amanda,", "fundadora da Recria, agência de marketing e comercial desde 2015."),
    ("Atende", "negócios B2C e B2B, de vários tamanhos e nichos: saúde, jurídico, varejo, indústria, serviços."),
    ("Já ajudou", "marcas a saírem de R$10 mil, R$80 mil/mês, para mais de R$1 milhão/mês de faturamento."),
], size=15, color=WHITE)
rect(s, 8.3, 2.1, 4.1, 4.4, fill=GOLD, radius=True)
textbox(s, 8.6, 2.4, 3.5, 0.5, "AO FINAL DESSA AULA", size=12, color=BLACK, bold=True, font=BODY_FONT)
textbox(s, 8.6, 2.95, 3.5, 3.3, "Você vai entender exatamente por que seu marketing não converte em venda e o que precisa existir no seu negócio para isso mudar.", size=14.5, color=BLACK, font=BODY_FONT, line_spacing=1.2)
pagenum(s, nextn())

# ============================================================
# SLIDE 3: COMO FUNCIONA ESSA AULA
# ============================================================
s = add_slide(BLACK)
kicker(s, "Como funciona essa aula")
textbox(s, 0.9, 1.1, 10.5, 1.6, "Isso é um aulão de marketing de verdade", size=34, color=WHITE, bold=True, font=TITLE_FONT, line_spacing=1.1)
textbox(s, 0.9, 2.55, 10.2, 1.3,
        "Atendi negócio local, ecommerce, consultoria, mentoria, SaaS, infoproduto, indústria, de quem está começando até quem já fatura milhões. O princípio é sempre o mesmo: você vai entender exatamente o que precisa existir no seu marketing e no seu comercial para vender mais.",
        size=16.5, color=GOLD_SOFT, font=BODY_FONT, line_spacing=1.3)
rect(s, 0.9, 4.4, 11.0, 1.6, fill=GOLD, radius=True)
textbox(s, 1.3, 4.9, 10.2, 0.7, "Fica até o final, tem um presente para você.", size=20, color=BLACK, bold=True, font=BODY_FONT)
pagenum(s, nextn())

# ============================================================
# BLOCO 1 — REDE SOCIAL E FUNIL
# ============================================================
block_title(1, "Toda empresa precisa de rede social hoje", "Não importa se você vende local, físico, online, ou os dois.")

s = add_slide(BLACK)
kicker(s, "Bloco 1 · Níveis de consciência")
textbox(s, 0.9, 1.0, 11, 0.8, "Cada pessoa que te vê está num momento diferente", size=27, color=WHITE, bold=True, font=TITLE_FONT)
niveis = [
    ("01", "Não sabe que tem o problema", 0.9),
    ("02", "Sabe do problema, não conhece a solução", 3.7),
    ("03", "Conhece a solução, não confia em você ainda", 6.5),
    ("04", "Confia em você, está pronta para comprar", 9.3),
]
for num, txt, x in niveis:
    rect(s, x, 2.2, 3.1, 4.2, fill=CARD_BG, radius=True)
    textbox(s, x + 0.3, 2.5, 2.5, 0.6, num, size=24, color=GOLD, bold=True, font=TITLE_FONT)
    textbox(s, x + 0.3, 3.3, 2.5, 2.9, txt, size=14, color=WHITE, font=BODY_FONT, line_spacing=1.25)
pagenum(s, nextn())

quote_slide("“Imagine sua loja numa rua escura, onde quase ninguém passa. O anúncio é levar essa loja para uma avenida movimentada, ou para dentro de um shopping cheio de gente interessada em comprar.”")

# ============================================================
# BLOCO 2 — O ANÚNCIO NÃO FAZ MILAGRE
# ============================================================
block_title(2, "O anúncio acelera, não faz milagre", "Se a casa não estiver arrumada, o anúncio só leva mais gente para ver a desorganização.")

# ============================================================
# BLOCO 3 — ARRUMANDO A CASA
# ============================================================
block_title(3, "Arrumando a casa", "Site institucional, página de captação, ou LP de vendas direto para o checkout.")

two_col_compare(
    "Qual presença digital faz sentido para você",
    "NEGÓCIO MISTO OU FÍSICO", [
        "Site institucional: mostra quem você é, sua história, gera confiança de marca",
        "Página de captação: formulário onde a pessoa deixa o contato",
        "Time comercial entra em contato depois",
    ],
    "NEGÓCIO 100% ONLINE", [
        "LP de vendas: explica tudo, funciona como um vendedor",
        "Gera atenção, desejo, oferta e fechamento",
        "Leva direto para o checkout (ex: Hotmart)",
    ],
    colA_gold=False, colB_gold=True,
)

# ============================================================
# BLOCO 4 — TREINAR O COMERCIAL
# ============================================================
block_title(4, "Treinar o comercial para agir", "A casa está arrumada. Agora alguém precisa saber atender quem chega.")

s = add_slide(BLACK)
kicker(s, "Bloco 4 · Velocidade e persistência")
textbox(s, 0.9, 1.0, 11, 0.8, "Atender rápido não é suficiente, atender persistente é o que fecha", size=24, color=WHITE, bold=True, font=TITLE_FONT, line_spacing=1.15)
rect(s, 0.9, 2.2, 5.6, 4.3, fill=CARD_BG, radius=True)
textbox(s, 1.25, 2.5, 5.0, 0.5, "VELOCIDADE", size=14, color=GOLD, bold=True, font=BODY_FONT)
bullets(s, 1.25, 3.05, 4.9, 3.2, [
    "Lead entrou, comercial recebe notificação e atende o mais rápido possível",
    "Ressalva combinada: lead de sexta à noite ou fim de semana, atende na segunda",
], size=14.5, color=WHITE)
rect(s, 6.85, 2.2, 5.6, 4.3, fill=GOLD, radius=True)
textbox(s, 7.2, 2.5, 5.0, 0.5, "PERSISTÊNCIA", size=14, color=BLACK, bold=True, font=BODY_FONT)
bullets(s, 7.2, 3.05, 4.9, 3.2, [
    "Mínimo 5 tentativas de contato antes de marcar como desqualificado",
    "Canais diferentes: WhatsApp, Instagram e e-mail",
    "Facilitar o pagamento: parcelado, Pix, boleto, cartão",
], size=14.5, color=BLACK, bullet_color=BLACK)
pagenum(s, nextn())

# ============================================================
# BLOCO 5 — CANAIS DE ANÚNCIO
# ============================================================
block_title(5, "Cada canal de anúncio funciona diferente", "Google é intenção. Instagram é atenção.")

two_col_compare(
    "Por que o custo do lead muda de canal para canal",
    "GOOGLE · INTENÇÃO", [
        "“Furei o pneu no meio da rua, vou direto no Google buscar borracharia”",
        "Lead mais perto da decisão, mais fundo de funil",
        "Por isso tende a ser mais caro",
    ],
    "INSTAGRAM · ATENÇÃO", [
        "Ninguém entra na rede social querendo comprar",
        "O anúncio rouba a atenção em um momento de distração",
        "Jornada mais longa, por isso o lead tende a ser mais barato",
    ],
    colA_gold=False, colB_gold=False,
)

# Slide em branco para Amanda inserir exemplos reais de anúncio (Google e Instagram)
s = add_slide(BLACK)
kicker(s, "Bloco 5 · Exemplos de anúncio")
textbox(s, 0.9, 3.4, 11.0, 1.0, "[Inserir exemplos de anúncios aqui]", size=20, color=GRAY, italic=True, font=BODY_FONT, align=PP_ALIGN.CENTER)
pagenum(s, nextn())

s = add_slide(BLACK)
kicker(s, "Bloco 5 · A jornada no Instagram")
textbox(s, 0.9, 1.0, 11, 0.8, "Nem todo mundo que vê o anúncio decide igual", size=28, color=WHITE, bold=True, font=TITLE_FONT)
caminhos = [
    ("Compra na hora", "decide e fecha no mesmo momento que vê o anúncio."),
    ("Se conecta, compra depois", "segue o perfil, acompanha um pouco, e fecha mais adiante."),
    ("Acompanha meses sem comprar", "fica de olho por um bom tempo, sem decidir, mas pode comprar um dia."),
    ("Ignora", "passa direto, não se conecta com a oferta."),
]
bullets(s, 0.9, 2.5, 10.8, 4.0, caminhos, size=17, color=WHITE, space_after=24, line_spacing=1.2)
pagenum(s, nextn())

# ============================================================
# BLOCO 6 — FLUXO COMPLETO NA PRÁTICA
# ============================================================
block_title(6, "O fluxo completo na prática", "Do anúncio até o follow-up proativo.")

s = add_slide(BLACK)
kicker(s, "Bloco 6 · O caminho do lead")
textbox(s, 0.9, 1.0, 11, 0.8, "Mesmo sem resposta, você não perde o contato", size=27, color=WHITE, bold=True, font=TITLE_FONT)
passos = ["Anúncio no Instagram", "Página com formulário", "Dados capturados", "Redireciona para o WhatsApp", "Follow-up proativo, 5+ tentativas"]
circle_flow(s, passos, centers_y=3.2, circ_d=0.85, spacing=2.45, startx=1.3, label_y=4.25,
            label_w=2.3, num_size=22, label_size=12.5, highlight_idx=4)
pagenum(s, nextn())

# Jornada linear x não linear (na prática)
s = add_slide(BLACK)
kicker(s, "Bloco 6 · Na prática")
textbox(s, 0.9, 1.0, 11, 0.8, "Nem toda jornada até a venda é igual", size=27, color=WHITE, bold=True, font=TITLE_FONT)
rect(s, 0.9, 2.05, 11.0, 2.15, fill=CARD_BG, radius=True)
textbox(s, 1.25, 2.3, 4.5, 0.4, "JORNADA LINEAR", size=13, color=GOLD, bold=True, font=BODY_FONT)
textbox(s, 1.25, 2.75, 10.3, 1.3,
        "Vê o anúncio → conhece na rede social → clica → preenche o formulário na página → vê o botão do WhatsApp → comercial atende (pelo WhatsApp ou pelos dados do formulário) → venda.",
        size=13.5, color=WHITE, font=BODY_FONT, line_spacing=1.3)
rect(s, 0.9, 4.4, 11.0, 2.5, fill=GOLD, radius=True)
textbox(s, 1.25, 4.65, 4.5, 0.4, "JORNADA NÃO LINEAR", size=13, color=BLACK, bold=True, font=BODY_FONT)
textbox(s, 1.25, 5.1, 10.3, 1.6,
        "Vê o anúncio, mas em vez de deixar o dado, segue a rede social. Pesquisa no Google, compara, some, volta semanas depois. Pode comprar de cara e voltar a comprar, ou comprar só depois de acompanhar por um tempo, e ainda indicar para outra pessoa.",
        size=13.5, color=BLACK, font=BODY_FONT, line_spacing=1.3)
pagenum(s, nextn())

quote_slide(
    "“É por isso que, às vezes, um anúncio demora um pouco mais ou um pouco menos para dar resultado.”"
)

# ============================================================
# BLOCO 7 — MARKETING NÃO VENDE, COMERCIAL VENDE
# ============================================================
block_title(7, "Marketing não vende, comercial vende", "Pare de culpar o marketing por uma venda perdida.")

quote_slide(
    "“O trabalho do marketing é levar o lead o mais próximo possível de estar interessado. Quem quebra a objeção e converte é o comercial.”",
    "Se o comercial não sabe fazer follow-up nem facilitar o pagamento, nem o melhor marketing do mundo resolve."
)

# ============================================================
# BLOCO 8 — AUTOMAÇÃO
# ============================================================
block_title(8, "Automação amarrando tudo", "Conecta as pontas sem depender de alguém checar planilha o dia inteiro.")

s = add_slide(BLACK)
kicker(s, "Bloco 8 · Temperatura do lead")
textbox(s, 0.9, 1.0, 10, 0.8, "Nem todo lead está no mesmo estágio", size=28, color=WHITE, bold=True, font=TITLE_FONT)
defs = [
    ("LEAD", "Alguém que demonstrou interesse e deixou um dado de contato.", 0.9),
    ("MQL", "Tomou uma ação real (baixou material, assistiu até o fim), mas ainda não validado pelo comercial.", 4.7),
    ("SQL", "Validado pelo comercial: tem perfil, orçamento e momento para comprar.", 8.5),
]
for label, body, x in defs:
    rect(s, x, 2.2, 3.7, 4.2, fill=CARD_BG, radius=True)
    textbox(s, x + 0.35, 2.55, 3.0, 0.6, label, size=22, color=GOLD, bold=True, font=TITLE_FONT)
    textbox(s, x + 0.35, 3.3, 3.0, 2.9, body, size=14, color=WHITE, font=BODY_FONT, line_spacing=1.25)
pagenum(s, nextn())

# ============================================================
# BLOCO 9 — CICLO DE VENDA POR TICKET
# ============================================================
block_title(9, "O ciclo de venda muda com o ticket", "Quanto mais caro o produto, mais esforço dos dois lados.")

s = add_slide(BLACK)
kicker(s, "Bloco 9 · Ticket baixo, médio e alto")
textbox(s, 0.9, 1.0, 11, 0.8, "O mesmo princípio em qualquer tipo de negócio", size=26, color=WHITE, bold=True, font=TITLE_FONT)
tickets = [
    ("TICKET BAIXO", [
        ("Infoproduto", "Ebook ou mini-curso de R$27 a R$97"),
        ("Produto físico", "Item de loja, oferta do dia"),
        ("Consultoria/serviço", "Sessão avulsa, ajuste pontual"),
    ], "Decisão quase instantânea, pouco esforço de marketing e comercial", False),
    ("TICKET MÉDIO", [
        ("Infoproduto", "Curso completo de R$300 a R$1.500"),
        ("Produto físico", "Equipamento, mobília, eletrônico"),
        ("Consultoria/serviço", "Pacote mensal, projeto fechado"),
    ], "Precisa de mais prova social e algumas etapas de confiança", False),
    ("TICKET ALTO", [
        ("Infoproduto", "Mentoria de R$5 mil ou mais"),
        ("Produto físico", "Imóvel, carro, maquinário"),
        ("Consultoria/serviço", "Consultoria estratégica de longo prazo"),
    ], "Mais esforço dos dois lados: jornada longa, prova social forte, comercial próximo", True),
]
xw, gap = 3.65, 0.2
for i, (label, items, nota, gold) in enumerate(tickets):
    x = 0.9 + i * (xw + gap)
    fill = GOLD if gold else CARD_BG
    color = BLACK if gold else WHITE
    rect(s, x, 2.1, xw, 4.6, fill=fill, radius=True)
    textbox(s, x + 0.25, 2.35, xw - 0.5, 0.4, label, size=14, color=BLACK if gold else GOLD, bold=True, font=BODY_FONT)
    bullets(s, x + 0.25, 2.85, xw - 0.5, 2.3, items, size=12, color=color, bullet_color=color, space_after=8)
    textbox(s, x + 0.25, 5.25, xw - 0.5, 1.3, nota, size=11.5, color=color, italic=True, font=BODY_FONT, line_spacing=1.2)
pagenum(s, nextn())

two_col_compare(
    "Dois exemplos, dois esforços diferentes",
    "TICKET BAIXO · CICLO CURTO", [
        "Exemplo: oferta de carne no mercado",
        "Vê o anúncio, entra no grupo de oferta no WhatsApp",
        "Ou só vê o anúncio e vai comprar na promoção",
        "Decisão quase instantânea",
    ],
    "TICKET ALTO · CICLO LONGO", [
        "Exemplo: mentoria de R$5 mil",
        "LP bem feita, comercial chama com atenção total",
        "Gera proximidade, autoridade, mostra prova social",
        "Facilita: parcelado, boleto, Pix, cartão",
    ],
    colA_gold=False, colB_gold=True,
)

# ============================================================
# TRANSIÇÃO — VALIDAÇÃO EMOCIONAL
# ============================================================
quote_slide(
    "“Se você nunca tinha pensado em nada disso antes, não se cobre por isso. Ninguém te ensinou. É exatamente por isso que você está aqui agora.”"
)

# ============================================================
# LINHA DO TEMPO / MAPA MENTAL
# ============================================================
s = add_slide(BLACK)
kicker(s, "Antes das provas")
textbox(s, 0.9, 1.0, 11, 0.8, "A linha do tempo de um negócio que funciona", size=27, color=WHITE, bold=True, font=TITLE_FONT)
etapas = [
    ("Ideia", "Negócio é idealizado"),
    ("Nasce", "Negócio é aberto"),
    ("Marketing entra", "Atrai atenção certa"),
    ("Comercial entra", "Conduz até o sim"),
    ("Lead → Cliente", "Marketing e comercial se conectam"),
    ("Cliente recorrente", "Compra de novo"),
    ("Indicação", "Cliente traz outro cliente"),
]
circle_flow(s, etapas, centers_y=2.95, circ_d=0.62, spacing=1.62, startx=1.35, label_y=3.55,
            label_w=1.6, num_size=16, label_size=11.5, gold_idxs=(2, 3))
pagenum(s, nextn())

s = add_slide(BLACK)
kicker(s, "Antes das provas · Dados de mercado")
textbox(s, 0.9, 1.0, 11, 0.8, "Quando marketing e comercial andam juntos", size=28, color=WHITE, bold=True, font=TITLE_FONT)
dados = [
    ("+20%", "de crescimento de receita em empresas com marketing e comercial alinhados"),
    ("+67%", "mais eficiência em fechar negócios quando os dois times trabalham integrados"),
    ("até 95%", "de aumento no lucro com apenas 5% de aumento na retenção de clientes"),
]
xs = [0.9, 4.85, 8.8]
for (num, desc), x in zip(dados, xs):
    rect(s, x, 2.3, 3.6, 3.9, fill=CARD_BG, radius=True)
    textbox(s, x + 0.3, 2.6, 3.0, 1.0, num, size=34, color=GOLD, bold=True, font=TITLE_FONT)
    textbox(s, x + 0.3, 3.7, 3.0, 2.3, desc, size=13.5, color=WHITE, font=BODY_FONT, line_spacing=1.25)
textbox(s, 0.9, 6.5, 11, 0.5, "Dados ilustrativos de estudos de mercado amplamente citados sobre alinhamento entre marketing e vendas (ex.: SiriusDecisions, Aberdeen, Bain & Co.) — confirme as fontes antes de citar números exatos publicamente.", size=10, color=GRAY, italic=True, font=BODY_FONT)
pagenum(s, nextn())

# ============================================================
# BLOCO 10 — PROVAS: CASES REAIS
# ============================================================
block_title(10, "Provas reais", "Dois negócios, dois pontos de partida, o mesmo princípio.")

s = add_slide(BLACK)
kicker(s, "Case 1 · Suplemento alimentar")
textbox(s, 0.9, 1.0, 11, 0.8, "Cliente que já existia, mas nunca foi ativado", size=26, color=WHITE, bold=True, font=TITLE_FONT)
rect(s, 0.9, 2.1, 6.7, 4.4, fill=CARD_BG, radius=True)
textbox(s, 1.25, 2.4, 6.0, 0.5, "ANTES", size=13, color=GOLD, bold=True, font=BODY_FONT)
textbox(s, 1.25, 2.9, 6.0, 3.4,
        "Empresa com quase 5 anos, sem CRM funcionando. Ninguém conseguia ativar a base de clientes.\n\nFoi feito: estudo fino do público, comunicação personalizada, WhatsApp e e-mail alternando conteúdo educativo e promoção.",
        size=14, color=WHITE, font=BODY_FONT, line_spacing=1.3)
rect(s, 7.85, 2.1, 4.55, 4.4, fill=GOLD, radius=True)
textbox(s, 8.2, 2.4, 3.9, 0.4, "RESULTADO", size=12, color=BLACK, bold=True, font=BODY_FONT)
textbox(s, 8.2, 2.9, 3.9, 0.7, "R$200 mil", size=25, color=BLACK, bold=True, font=TITLE_FONT)
textbox(s, 8.2, 3.55, 3.9, 0.4, "1º mês", size=12, color=BLACK, font=BODY_FONT)
textbox(s, 8.2, 4.0, 3.9, 0.7, "R$400 mil", size=25, color=BLACK, bold=True, font=TITLE_FONT)
textbox(s, 8.2, 4.65, 3.9, 0.4, "2º mês", size=12, color=BLACK, font=BODY_FONT)
textbox(s, 8.2, 5.1, 3.9, 0.7, "R$500 mil", size=25, color=BLACK, bold=True, font=TITLE_FONT)
textbox(s, 8.2, 5.75, 3.9, 0.4, "3º mês · sem tráfego pago", size=12, color=BLACK, font=BODY_FONT)
pagenum(s, nextn())

s = add_slide(BLACK)
kicker(s, "Case 2 · Portas de enrolar")
textbox(s, 0.9, 1.0, 11, 0.8, "De negócio local a expansão regional", size=27, color=WHITE, bold=True, font=TITLE_FONT)
rect(s, 0.9, 2.1, 6.7, 4.4, fill=CARD_BG, radius=True)
textbox(s, 1.25, 2.4, 6.0, 0.5, "O QUE FOI FEITO", size=13, color=GOLD, bold=True, font=BODY_FONT)
bullets(s, 1.25, 2.95, 6.0, 3.4, [
    "Páginas diferentes para cada público (cliente final x intermediador)",
    "Clube de parceria e indicação",
    "Ajuste de anúncios em Google e Meta",
    "Resposta comercial em até 1 dia",
], size=13.5, color=WHITE)
rect(s, 7.85, 2.1, 4.55, 4.4, fill=GOLD, radius=True)
textbox(s, 8.2, 2.4, 3.9, 0.4, "EVOLUÇÃO", size=12, color=BLACK, bold=True, font=BODY_FONT)
textbox(s, 8.2, 2.85, 3.9, 0.5, "R$80 mil/mês → R$500 mil/mês em 4 meses", size=13, color=BLACK, bold=True, font=BODY_FONT, line_spacing=1.2)
textbox(s, 8.2, 3.55, 3.9, 0.5, "1º milhão em 10 meses", size=13, color=BLACK, bold=True, font=BODY_FONT)
textbox(s, 8.2, 4.15, 3.9, 0.5, "R$1,3 milhão aos 12 meses", size=13, color=BLACK, bold=True, font=BODY_FONT)
textbox(s, 8.2, 4.75, 3.9, 0.7, "2º milhão em 1 ano e meio. Hoje expandindo para Santa Catarina.", size=12.5, color=BLACK, font=BODY_FONT, line_spacing=1.2)
pagenum(s, nextn())

quote_slide(
    "“Todo negócio, sem exceção, precisa de marketing bom e comercial afiado para crescer de verdade.”",
    "Hotel, loja, ecommerce, salão, curso, consultoria, universidade, mentoria, até curso de meditação para dormir. O foco muda, o princípio não."
)

# ============================================================
# BLOCO 10.5 — O RESUMÃO
# ============================================================
s = add_slide(BLACK)
kicker(s, "O resumão")
textbox(s, 0.9, 1.0, 11, 0.8, "Tudo que você viu até aqui", size=28, color=WHITE, bold=True, font=TITLE_FONT)
resumo_itens = [
    "Rede social com funil aplicado (consciência)",
    "Anúncio acelera, não faz milagre",
    "Casa arrumada: site institucional ou LP",
    "Comercial treinado: velocidade, follow-up, pagamento facilitado",
    "Canal certo para o seu negócio: Google x Instagram",
    "Fluxo completo amarrado, do anúncio ao follow-up",
    "Marketing leva até a venda, comercial fecha a venda",
    "Automação conectando tudo",
    "Esforço proporcional ao ticket do seu produto",
]
col1 = resumo_itens[:5]
col2 = resumo_itens[5:]
bullets(s, 0.9, 2.1, 5.7, 4.5, col1, size=14.5, color=WHITE, space_after=14)
bullets(s, 6.85, 2.1, 5.6, 4.5, col2, size=14.5, color=WHITE, space_after=14)
rect(s, 0.9, 6.55, 11.5, 0.7, fill=GOLD, radius=True)
textbox(s, 1.2, 6.72, 11.0, 0.4, "Baixe o material de apoio com esse resumo abaixo do vídeo.", size=14, color=BLACK, bold=True, font=BODY_FONT)
pagenum(s, nextn())

# ============================================================
# BLOCO 11 — O PRESENTE E A OFERTA
# ============================================================
s = add_slide(BLACK)
kicker(s, "Você chegou até aqui")
textbox(s, 0.9, 1.6, 11.0, 1.3, "Parabéns por chegar até aqui! Aqui está o seu primeiro presente!", size=28, color=WHITE, bold=True, font=TITLE_FONT, line_spacing=1.15)
textbox(s, 0.9, 3.1, 10.6, 1.6,
        "Nessa mesma página tem um botão para você baixar, totalmente grátis, o material de apoio em PDF: para lembrar de tudo e já começar a aplicar no dia a dia, conferindo como está o seu negócio.",
        size=16, color=GOLD_SOFT, font=BODY_FONT, line_spacing=1.3)
textbox(s, 0.9, 4.85, 10.6, 1.3,
        "E se você quiser o meu apoio, minha ajuda para estruturar ou reestruturar o seu marketing e o seu comercial, tenho mais um presente para você.",
        size=16, color=WHITE, bold=True, font=BODY_FONT, line_spacing=1.3)
pagenum(s, nextn())

s = add_slide(BLACK)
kicker(s, "Bloco 11 · Seu presente")
textbox(s, 0.9, 1.05, 11, 1.0, "Um diagnóstico gratuito, dois caminhos", size=28, color=WHITE, bold=True, font=TITLE_FONT, line_spacing=1.1)
rect(s, 0.9, 2.3, 11.0, 3.9, fill=CARD_BG, radius=True)
textbox(s, 1.3, 2.65, 5.0, 0.5, "DIAGNÓSTICO GRATUITO", size=14, color=GOLD, bold=True, font=BODY_FONT)
textbox(s, 1.3, 3.2, 10.2, 2.7,
        "Você pode começar já com um diagnóstico gratuito seguido do Pacote Inicial (R$1.500/mês), onde eu entrego tudo que você precisa para começar agora no seu negócio, condição válida só nessa página. Ou pode aplicar para uma Sessão Estratégica comigo: eu entendo o cenário do seu negócio e te entrego o diagnóstico de qualquer maneira, você decide depois se faz sentido fechar ou não com a minha agência.",
        size=15, color=WHITE, font=BODY_FONT, line_spacing=1.35)
pagenum(s, nextn())

s = add_slide(BLACK)
kicker(s, "Bloco 11 · Seus dois caminhos")
textbox(s, 0.9, 1.0, 10, 0.8, "Dois convites para você agora", size=28, color=WHITE, bold=True, font=TITLE_FONT)
rect(s, 0.9, 2.1, 5.6, 4.5, fill=GOLD, radius=True)
textbox(s, 1.25, 2.4, 5.0, 0.4, "CAMINHO MAIS RÁPIDO", size=12, color=BLACK, bold=True, font=BODY_FONT)
textbox(s, 1.25, 2.8, 5.0, 0.5, "PACOTE INICIAL", size=14, color=BLACK, bold=True, font=BODY_FONT)
textbox(s, 1.25, 3.25, 5.0, 0.6, "R$1.500/mês", size=22, color=BLACK, bold=True, font=TITLE_FONT)
bullets(s, 1.25, 3.95, 4.9, 2.5, [
    ("Tráfego pago:", "anúncios no Google e no Instagram"),
    ("Social media:", "criação de conteúdo nas redes sociais"),
    ("Site:", "captação de potenciais clientes com automação que avisa o time comercial"),
    ("Treinamento comercial:", "para o seu time não perder mais venda"),
], size=12, color=BLACK, bullet_color=BLACK, space_after=7)
rect(s, 6.85, 2.1, 5.6, 4.5, fill=CARD_BG, radius=True)
textbox(s, 7.2, 2.4, 5.0, 0.4, "CAMINHO MAIS DEVAGAR, COM CAUTELA", size=12, color=GOLD, bold=True, font=BODY_FONT)
textbox(s, 7.2, 2.8, 5.0, 0.5, "SESSÃO ESTRATÉGICA", size=14, color=WHITE, bold=True, font=BODY_FONT)
textbox(s, 7.2, 3.25, 5.0, 0.7, "Aplicação, vagas limitadas", size=18, color=WHITE, bold=True, font=TITLE_FONT)
textbox(s, 7.2, 4.0, 4.9, 2.5,
        "Não é uma reunião de vendas, é um diagnóstico que eu mesma farei para o seu negócio. O diagnóstico é seu de qualquer forma, fechando ou não algum serviço com a minha agência. Você sai sabendo onde está, para onde quer ir e quais são os gargalos. Você só tem a ganhar.",
        size=12.5, color=WHITE, font=BODY_FONT, line_spacing=1.28)
pagenum(s, nextn())

s = add_slide(BLACK)
kicker(s, "Bloco 11 · O caminho que custa caro")
textbox(s, 0.9, 1.1, 11, 1.2, "Ou você pode continuar tentando sozinho", size=27, color=WHITE, bold=True, font=TITLE_FONT, line_spacing=1.1)
rect(s, 0.9, 2.6, 11.0, 3.7, fill=CARD_BG, radius=True)
bullets(s, 1.3, 3.0, 10.2, 3.0, [
    "Gastando sem direção",
    "Perdendo venda para concorrente",
    "Correndo o risco de chegar até a falência",
], size=18, color=WHITE, space_after=16)
textbox(s, 1.3, 5.6, 10.0, 0.6, "Negócio nenhum cresce sem marketing e comercial alinhados, independente do que você vende.", size=15, color=GOLD_SOFT, italic=True, font=BODY_FONT, line_spacing=1.25)
pagenum(s, nextn())

s = add_slide(BLACK)
logo_mark(s, x=0.9, y=0.6)
textbox(s, 0.9, 2.1, 11.0, 1.2, "Qual caminho você escolhe agora?", size=36, color=WHITE, bold=True, font=TITLE_FONT, align=PP_ALIGN.CENTER)
btn_w, btn_h, btn_y = 5.2, 1.7, 3.8
rect(s, 0.9, btn_y, btn_w, btn_h, fill=GOLD, radius=True)
textbox(s, 1.1, btn_y, btn_w - 0.4, btn_h, "Quero o\nPacote Inicial", size=19, color=BLACK, bold=True,
        font=BODY_FONT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.2)
rect(s, 6.3, btn_y, btn_w, btn_h, fill=CARD_BG, radius=True, line=GOLD)
textbox(s, 6.5, btn_y, btn_w - 0.4, btn_h, "Quero aplicar para a\nSessão Estratégica", size=17, color=GOLD, bold=True,
        font=BODY_FONT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.2)
textbox(s, 0.9, 6.0, 11.0, 0.5, "Os botões estão na página abaixo do vídeo.", size=13, color=GRAY, italic=True, font=BODY_FONT, align=PP_ALIGN.CENTER)
pagenum(s, nextn())

# ============================================================
# SLIDE FINAL — OBRIGADA
# ============================================================
s = add_slide(BLACK)
logo_mark(s, x=0.9, y=0.6)
textbox(s, 0.9, 3.0, 11.5, 1.5, "Obrigada!", size=48, color=WHITE, bold=True, font=TITLE_FONT, align=PP_ALIGN.CENTER)
textbox(s, 0.9, 4.4, 11.5, 0.7, "por chegar até aqui comigo.", size=18, color=GOLD_SOFT, italic=True, font=BODY_FONT, align=PP_ALIGN.CENTER)

prs.save("aulao-recria-slides.pptx")
print("OK", n, "slides")
